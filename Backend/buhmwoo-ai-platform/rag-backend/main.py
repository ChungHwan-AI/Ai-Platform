# main.py

import logging
import os
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import List, Optional, Dict, Any

import httpx
from dotenv import load_dotenv
from fastapi import FastAPI, UploadFile, File, HTTPException, Form, Request
from fastapi.responses import HTMLResponse
from fastapi.templating import Jinja2Templates
from langchain_core.documents import Document
from openai import OpenAI
from pydantic import (
    BaseModel,
    Field,
    ConfigDict,
    AliasChoices,
    model_validator,
)

from config_chroma import get_vectordb, get_chroma_settings, reset_collection
from config_embed import get_embedding_fn, get_embedding_backend_info_dict
from utils.encoding import to_utf8_text
from utils.chunking import (
    chunk_single_document,
    DEFAULT_CHUNK_OVERLAP,
    DEFAULT_CHUNK_SIZE,
)
from admin_router import router as admin_router
from rag_service import query_text

# -----------------------------
# 로깅 설정
# -----------------------------

logger = logging.getLogger(__name__)

if not logging.getLogger().handlers:
    logging.basicConfig(level=logging.INFO)

logging.getLogger("uvicorn").setLevel(logging.INFO)

# -----------------------------
# 환경 변수 로드
# -----------------------------

load_dotenv()

# -----------------------------
# OpenAI Responses + web_search 설정
# -----------------------------

openai_client = OpenAI(
    api_key=os.environ.get("OPENAI_API_KEY"),
    base_url=os.environ.get("OPENAI_BASE_URL") or None,
)

OPENAI_RESPONSES_MODEL = (
    os.environ.get("OPENAI_RESPONSES_MODEL")
    or os.environ.get("OPENAI_CHAT_MODEL")
    or "gpt-4.1-mini"
)

# -----------------------------
# 프롬프트 관련 구조체/문자열
# -----------------------------


@dataclass
class PromptContextChunk:
    """프롬프트 및 응답 모두에서 재사용할 검색 청크 정보를 표현하는 구조체"""

    reference: str  # LLM 응답 인용에 사용될 라벨 (예: [청크 1])
    header: str  # 프롬프트에 노출될 출처 요약 헤더
    body: str  # 청크 본문 텍스트
    metadata: Dict[str, Any]  # 원본 문서 메타데이터 사본
    source: str  # UX에서 표시할 주요 출처(파일명 또는 문서 ID)
    page: Optional[int]  # 페이지/슬라이드 등 위치 정보
    chunk_index: int  # 1부터 시작하는 청크 순번
    preview: str  # 프론트에서 사용할 간단한 스니펫 텍스트


PROMPT_SYSTEM_TEXT = (
    "당신은 한국어로 답변하는 친절한 업무 도우미입니다."
    "\n- [참고 문서 및 웹 검색 컨텍스트] 블록에 포함된 정보를 우선 활용해 질문에 답변하세요."
    "\n- 웹 검색 결과가 문서 내용이나 당신의 기존 지식과 다를 경우, 최신 외부 정보(웹 검색 결과)를 더 신뢰하세요."
    "\n- 먼저 사용자의 질문에 대한 결론을 1~2문장으로 분명하게 말한 뒤, 필요하면 간단한 설명을 1~3문장 덧붙이세요."
    "\n- 너무 딱딱한 면책 문구(예: '근거 부족', '정책과 다를 수 있습니다')는 사용하지 말고,"
    " 필요할 때만 짧게 '정확한 최신 정보는 공식 자료에서 한 번 더 확인해 주세요.' 정도만 덧붙이세요."
    "\n- 답변은 자연스러운 한국어 문단 또는 짧은 목록 형태로 작성하고,"
    " '## 핵심 요약', '## 사실에 기반한 창의적 통찰' 같은 섹션 헤더는 사용하지 마세요."
)

PROMPT_HUMAN_TEMPLATE = (
    "다음은 사용자의 질문과 참고용으로 제공된 문서/검색 컨텍스트입니다.\n\n"
    "[질문]\n"
    "{question}\n\n"
    "[참고 문서 및 웹 검색 컨텍스트]\n"
    "{context}\n\n"
    "위 정보를 바탕으로, 사용자가 실제 업무에 바로 활용할 수 있을 정도로 "
    "간단명료하고 자연스러운 한국어로 답변해 주세요."
)

PROMPT_SYSTEM_GENERAL = (
    "당신은 한국어로 답변하는 친절한 업무 비서입니다."
    "\n- 회사 문서가 없더라도 일반 상식과 논리를 활용해 문제를 해결하세요."
    "\n- 사실 근거가 부족하면 추측임을 명확히 밝히고, 안전하고 책임감 있게 답하세요."
    "\n- 사용자가 요청한 작업(요약, 번역, 일정 제안 등)을 그대로 수행하세요."
)

PROMPT_HUMAN_GENERAL = (
    "사용자 요청: {question}\n"
    "참고 문서: {context}\n"
    "친절하게 한국어로 응답하세요."
)


def _build_openai_prompt(question: str, context: Optional[str], has_context: bool) -> str:
    """
    OpenAI Responses API에 그대로 넘길 하나의 텍스트 프롬프트를 만든다.
    """
    q = (question or "").strip()
    c = (context or "").strip()

    if has_context and c:
        # 문서 기반 RAG 답변용
        system = PROMPT_SYSTEM_TEXT
        human = PROMPT_HUMAN_TEMPLATE.format(question=q, context=c)
    else:
        # 일반 지식 + 웹 검색용
        system = PROMPT_SYSTEM_GENERAL
        human = PROMPT_HUMAN_GENERAL.format(question=q, context=c)

    return f"{system.strip()}\n\n{human.strip()}"


# -----------------------------
# Pydantic 모델
# -----------------------------


class QueryRequest(BaseModel):
    """/query 요청 스키마 정의"""

    question: str = Field(..., description="사용자가 던진 질문 문장")
    doc_id: Optional[str] = Field(
        default=None,
        description="벡터 검색 대상을 제한할 문서 ID (선택)",
        alias="docId",
        validation_alias=AliasChoices("docId", "doc_id"),
    )
    top_k: int = Field(
        default=4,
        ge=1,
        le=20,
        description="벡터 검색으로 가져올 상위 청크 개수",
    )

    model_config = ConfigDict(
        populate_by_name=True,
        json_schema_extra={
            "examples": [
                {
                    "question": "이번 분기 생산 목표는 어떻게 되나요?",
                    "docId": "123e4567-e89b-12d3-a456-426614174000",
                    "top_k": 3,
                }
            ]
        },
    )


class Match(BaseModel):
    """벡터 검색 결과 청크 정보를 담는 스키마"""

    reference: str = Field(..., description="LLM 응답 인용에 사용할 청크 라벨")
    chunk_index: int = Field(
        ...,
        ge=1,
        alias="chunkIndex",
        validation_alias=AliasChoices("chunkIndex", "chunk_index"),
        description="검색 결과 순번 (1부터 시작)",
    )
    content: str = Field(..., description="검색된 청크 본문")
    preview: str = Field(..., description="줄바꿈 제거/길이 제한 미리보기 텍스트")
    source: str = Field(..., description="원본 문서 이름 또는 ID")
    page: Optional[int] = Field(
        default=None,
        description="문서 내 위치 정보 (해당하지 않으면 null)",
    )
    metadata: Dict[str, Any] = Field(
        default_factory=dict,
        description="청크에 연결된 원본 메타데이터",
    )

    model_config = ConfigDict(populate_by_name=True)


class QueryResponse(BaseModel):
    """/query 응답 스키마 정의"""

    answer: str = Field(..., description="LLM이 생성한 최종 응답")
    matches: List[Match] = Field(
        default_factory=list,
        description="벡터 검색으로 선택된 청크 목록",
    )


class RetrieveResponse(BaseModel):
    """/query/retrieve 응답 스키마 정의"""

    matches: List[Match] = Field(
        default_factory=list,
        description="벡터 검색으로 선택된 청크 목록",
    )
    context: str = Field(..., description="GPT 프롬프트에 사용할 컨텍스트 전체 텍스트")


class GenerateRequest(BaseModel):
    """/query/generate 요청 스키마 정의"""

    question: str = Field(..., description="사용자가 던진 질문 문장")
    context: Optional[str] = Field(
        default=None,
        description="벡터 검색 결과를 하나의 텍스트로 결합한 컨텍스트",
    )


class GenerateResponse(BaseModel):
    """/query/generate 응답 스키마 정의"""

    answer: str = Field(..., description="LLM이 생성한 최종 응답")


class RagHealthResponse(BaseModel):
    """RAG 백엔드의 기본 구성 상태를 한눈에 보여주는 헬스체크 응답"""

    ready: bool = Field(..., description="핵심 의존성이 모두 준비되었는지 여부")
    llm_provider: str = Field(..., alias="llmProvider", description="선택된 LLM 공급자 이름")
    llm_ready: bool = Field(..., alias="llmReady", description="LLM 호출 준비 여부")
    llm_error: Optional[str] = Field(
        default=None, alias="llmError", description="LLM 설정 문제 발생 시 사유"
    )
    embedding: Dict[str, Any] = Field(
        default_factory=dict, description="임베딩 백엔드 상태 요약"
    )
    chroma: Dict[str, Any] = Field(
        default_factory=dict, description="Chroma 영속 디렉터리 및 컬렉션 상태"
    )


class DocDeleteRequest(BaseModel):
    """벡터 DB에서 삭제할 문서 조건을 표현하는 모델"""

    doc_id: Optional[str] = Field(
        default=None,
        alias="docId",
        description="삭제할 문서의 UUID",
        validation_alias=AliasChoices("docId", "doc_id"),
    )
    source: Optional[str] = Field(
        default=None,
        description="삭제할 문서를 찾기 위한 파일명",
    )

    model_config = ConfigDict(populate_by_name=True)

    @model_validator(mode="after")
    def validate_identifier(cls, values: "DocDeleteRequest") -> "DocDeleteRequest":
        if not (values.doc_id or values.source):
            raise ValueError("docId 또는 source 중 하나는 반드시 전달되어야 합니다.")
        return values


# -----------------------------
# RAG 컨텍스트 준비 / 검색 복구
# -----------------------------


def _prepare_prompt_context(docs: List[Document]) -> tuple[List[PromptContextChunk], str, bool]:
    """검색된 청크를 프롬프트와 API 응답 모두에서 활용 가능한 형태로 가공"""

    context_chunks: List[PromptContextChunk] = []
    context_blocks: List[str] = []
    has_docs = False

    for idx, doc in enumerate(docs, start=1):
        has_docs = True

        meta = dict(doc.metadata or {})
        source_raw = meta.get("source") or meta.get("docId") or "unknown"
        source = str(source_raw)

        page_value = meta.get("page")
        if isinstance(page_value, (int, float)):
            page = int(page_value)
        elif isinstance(page_value, str) and page_value.isdigit():
            page = int(page_value)
        else:
            page = None

        reference = f"[청크 {idx}]"
        header = f"{reference} 출처: {source}"
        if page is not None:
            header += f" | 페이지: {page}"

        body = doc.page_content.strip()
        snippet = " ".join(body.split())
        if len(snippet) > 160:
            snippet = snippet[:157] + "..."

        context_blocks.append(f"{header}\n{body}")
        context_chunks.append(
            PromptContextChunk(
                reference=reference,
                header=header,
                body=body,
                metadata=meta,
                source=source,
                page=page,
                chunk_index=idx,
                preview=snippet,
            )
        )

    if not context_blocks:
        context_blocks.append("(참고할 문서를 찾지 못했습니다.)")

    context_text = "\n\n".join(context_blocks)
    return context_chunks, context_text, has_docs


def _run_retrieval_with_recovery(
    *,
    question: str,
    top_k: int,
    metadata_filter: Dict[str, Any] | None,
):
    """retriever 실행 시 임베딩 차원 불일치 등을 감지해 자동 복구."""

    try:
        return query_text(question, k=top_k, metadata_filter=metadata_filter)
    except ValueError as err:
        message = str(err)
        if "dimensionality" in message:
            logger.warning(
                "임베딩 차원 불일치 감지 → Chroma 컬렉션 초기화 후 재검색을 시도합니다."
            )
            try:
                reset_collection()
                return query_text(question, k=top_k, metadata_filter=metadata_filter)
            except Exception as retry_exc:
                logger.exception("컬렉션 초기화 후 재검색에 실패했습니다.")
                raise HTTPException(
                    status_code=500,
                    detail=(
                        "벡터 컬렉션 임베딩 차원이 맞지 않아 재시도에 실패했습니다. "
                        "`python wipe_chroma.py` 실행 후 문서를 다시 업로드하세요."
                    ),
                ) from retry_exc
        raise


# -----------------------------
# OpenAI Responses + web_search 호출
# -----------------------------


def _generate_answer(
    question: str,
    context: Optional[str],
    has_context: bool,
) -> str:
    """
    문서 컨텍스트 유무에 따라 프롬프트를 만들고
    OpenAI Responses API + web_search 도구로 최종 답변을 생성한다.
    """

    provider = os.getenv("LLM_PROVIDER", "openai").lower()
    if provider != "openai":
        raise HTTPException(
            status_code=500,
            detail=f"지원하지 않는 LLM_PROVIDER: {provider} (현재는 openai + web_search만 지원)",
        )

    if not question or not question.strip():
        raise HTTPException(
            status_code=400,
            detail="질문이 비어 있습니다.",
        )

    prompt = _build_openai_prompt(question, context, has_context)

    try:
        response = openai_client.responses.create(
            model=OPENAI_RESPONSES_MODEL,
            input=prompt,
            tools=[{"type": "web_search"}],
        )
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"LLM 호출 실패: {e}",
        ) from e

    text = getattr(response, "output_text", None)

    if not text:
        try:
            output = getattr(response, "output", None)
            if output and len(output) > 0:
                first = output[0]
                content_list = getattr(first, "content", None)
                if content_list:
                    first_content = content_list[0]
                    text = getattr(first_content, "text", None)
        except Exception:
            text = None

    if not text:
        raise HTTPException(
            status_code=500,
            detail="LLM 응답에서 텍스트를 추출하지 못했습니다.",
        )

    return text.strip()


# -----------------------------
# 파일별 텍스트 추출 유틸
# -----------------------------


def _extract_text_txt_like_bytes(content: bytes) -> str:
    return to_utf8_text(content)


def _extract_text_pdf(path: Path) -> str:
    import fitz  # PyMuPDF

    out = []
    with fitz.open(str(path)) as pdf:
        for page in pdf:
            out.append(page.get_text("text"))
    return "\n".join(out)


def _iter_block_items(parent):
    """docx 문서/셀에서 단락과 표를 순서대로 순회하기 위한 유틸 함수"""

    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import _Cell, Table
    from docx.text.paragraph import Paragraph

    parent_elm = parent._tc if isinstance(parent, _Cell) else parent.element.body
    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)


def _collect_table_text(table, out_lines: List[str]) -> None:
    """docx 표를 순회하며 셀 텍스트를 라인으로 변환해 누적"""

    from docx.text.paragraph import Paragraph
    from docx.table import Table

    for row in table.rows:
        row_texts: List[str] = []
        for cell in row.cells:
            cell_fragments: List[str] = []
            for item in _iter_block_items(cell):
                if isinstance(item, Paragraph):
                    text = item.text.strip()
                    if text:
                        cell_fragments.append(text)
                elif isinstance(item, Table):
                    _collect_table_text(item, out_lines)
            cell_text = " ".join(cell_fragments).strip()
            if cell_text:
                row_texts.append(cell_text)
        if row_texts:
            out_lines.append(" | ".join(row_texts))


def _extract_text_docx(path: Path) -> str:
    from docx import Document as Docx
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Docx(str(path))
    lines: List[str] = []

    for block in _iter_block_items(doc):
        if isinstance(block, Paragraph):
            text = block.text.strip()
            if text:
                lines.append(text)
        elif isinstance(block, Table):
            _collect_table_text(block, lines)

    return "\n".join(lines)


def _shape_texts(shape) -> List[str]:
    """python-pptx shape 에서 텍스트/테이블 텍스트 추출"""
    texts: List[str] = []

    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        try:
            for para in shape.text_frame.paragraphs:
                texts.append("".join(run.text for run in para.runs) or (para.text or ""))
        except Exception:
            pass

    if hasattr(shape, "has_table") and shape.has_table:
        try:
            tbl = shape.table
            for r in tbl.rows:
                row_vals = []
                for c in r.cells:
                    s = (c.text or "").replace("\r", " ").replace("\n", " ").strip()
                    if s:
                        row_vals.append(s)
                if row_vals:
                    texts.append(" | ".join(row_vals))
        except Exception:
            pass

    if hasattr(shape, "shapes"):
        try:
            for s in shape.shapes:
                texts.extend(_shape_texts(s))
        except Exception:
            pass

    return [t for t in texts if t and t.strip()]


def _extract_text_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    out = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_buf = [f"[Slide {idx}]"]
        for shape in slide.shapes:
            slide_buf.extend(_shape_texts(shape))
        try:
            if (
                slide.has_notes_slide
                and slide.notes_slide
                and slide.notes_slide.notes_text_frame
            ):
                notes_texts = [
                    p.text
                    for p in slide.notes_slide.notes_text_frame.paragraphs
                    if p.text and p.text.strip()
                ]
                if notes_texts:
                    slide_buf.append("[Notes]")
                    slide_buf.append("\n".join(notes_texts))
        except Exception:
            pass
        out.append("\n".join(slide_buf))
    return "\n\n".join(out)


def _extract_text_xlsx(path: Path, max_rows_per_sheet: int = 2000) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(filename=str(path), read_only=True, data_only=True)
    out: List[str] = []
    for ws in wb.worksheets:
        out.append(f"[Sheet] {ws.title}")
        rows = 0
        for row in ws.iter_rows(values_only=True):
            if rows >= max_rows_per_sheet:
                out.append("... (truncated)")
                break
            vals = []
            for v in row:
                if v is None:
                    continue
                s = str(v).strip()
                if s:
                    vals.append(s)
            if vals:
                out.append(" | ".join(vals))
            rows += 1
    return "\n".join(out)


SUPPORTED_EXTS = {".txt", ".csv", ".log", ".md", ".pdf", ".docx", ".pptx", ".xlsx"}


def extract_text_by_ext(dst: Path, raw_content: bytes) -> str:
    ext = dst.suffix.lower()
    if ext in {".txt", ".csv", ".log", ".md"}:
        return _extract_text_txt_like_bytes(raw_content)
    if ext == ".pdf":
        return _extract_text_pdf(dst)
    if ext == ".docx":
        return _extract_text_docx(dst)
    if ext == ".pptx":
        return _extract_text_pptx(dst)
    if ext == ".xlsx":
        return _extract_text_xlsx(dst)
    raise HTTPException(status_code=400, detail=f"지원하지 않는 확장자: {ext}")


# -----------------------------
# FastAPI 앱 설정
# -----------------------------

BASE_DIR = Path(__file__).parent

app = FastAPI()
app.include_router(admin_router)

templates = Jinja2Templates(directory=str(BASE_DIR / "templates"))


def _check_llm_status() -> tuple[bool, Optional[str]]:
    provider = os.getenv("LLM_PROVIDER", "openai").lower()
    if provider != "openai":
        return False, f"지원하지 않는 LLM_PROVIDER: {provider}"

    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        return False, "OPENAI_API_KEY가 설정되어 있지 않습니다."

    return True, None


@app.get(
    "/health/rag",
    response_model=RagHealthResponse,
    summary="RAG 필수 설정 헬스체크",
    tags=["health"],
)
async def rag_health() -> RagHealthResponse:
    llm_ready, llm_error = _check_llm_status()

    try:
        embedding_info = get_embedding_backend_info_dict(force_refresh=True)
    except Exception as exc:
        embedding_info = {"error": str(exc), "resolved_backend": None}

    chroma_dir, collection_name = get_chroma_settings()
    chroma_summary: Dict[str, Any] = {
        "persistPath": chroma_dir,
        "collection": collection_name,
    }
    try:
        vectordb = get_vectordb(get_embedding_fn())
        chroma_summary["count"] = vectordb._collection.count()
    except Exception as exc:
        chroma_summary["error"] = str(exc)

    ready = llm_ready and not embedding_info.get("error") and "error" not in chroma_summary

    return RagHealthResponse(
        ready=ready,
        llm_provider=os.getenv("LLM_PROVIDER", "openai"),
        llm_ready=llm_ready,
        llm_error=llm_error,
        embedding=embedding_info,
        chroma=chroma_summary,
    )


@app.get("/", response_class=HTMLResponse, summary="웹 챗 인터페이스", tags=["ui"])
async def chat_ui(request: Request) -> HTMLResponse:
    return templates.TemplateResponse("chat.html", {"request": request})


@app.post("/upload")
async def upload(
    file: UploadFile = File(...),
    doc_id: Optional[str] = Form(default=None, alias="docId"),
):
    try:
        uploads_dir = (Path(__file__).parent / "uploads").resolve()
        uploads_dir.mkdir(parents=True, exist_ok=True)
        dst = uploads_dir / file.filename

        content: bytes = await file.read()
        max_bytes = 50 * 1024 * 1024
        if len(content) > max_bytes:
            raise HTTPException(status_code=413, detail=f"파일이 너무 큽니다(최대 50MB). size={len(content)}")

        dst.write_bytes(content)

        ext = dst.suffix.lower()
        if ext not in SUPPORTED_EXTS:
            raise HTTPException(400, f"지원하지 않는 확장자: {ext}")
        text = extract_text_by_ext(dst, content)
        if not text.strip():
            raise HTTPException(400, "본문이 비어 있습니다.")

        chunk_size = DEFAULT_CHUNK_SIZE
        chunk_overlap = DEFAULT_CHUNK_OVERLAP
        docs = chunk_single_document(
            text,
            {
                "source": file.filename,
                "doctype": ext[1:],
                **({"docId": doc_id} if doc_id else {}),
            },
            chunk_size=chunk_size,
            overlap=chunk_overlap,
        )

        emb = get_embedding_fn()
        vectordb = get_vectordb(emb)

        deletion_filter = {"docId": doc_id} if doc_id else {"source": file.filename}
        deleted_chunks = 0
        try:
            existing = vectordb._collection.get(
                where=deletion_filter,
                include=[],
            )
            existing_ids = existing.get("ids") if isinstance(existing, dict) else None
            matched_ids = existing_ids or []

            if matched_ids:
                delete_result = vectordb._collection.delete(where=deletion_filter)
                if isinstance(delete_result, dict):
                    deleted_chunks = len(delete_result.get("ids") or [])
                elif isinstance(delete_result, (list, tuple, set)):
                    deleted_chunks = len(delete_result)
                elif isinstance(delete_result, int):
                    deleted_chunks = delete_result
            logger.info(
                "업로드 전 기존 청크 삭제 수행 filter=%s, deleted=%s",
                deletion_filter,
                deleted_chunks,
            )
        except Exception as delete_exc:
            logger.exception("기존 청크 삭제 중 오류 발생 filter=%s", deletion_filter)
            raise HTTPException(
                status_code=500,
                detail="기존 문서를 삭제하지 못해 업로드를 중단합니다.",
            ) from delete_exc

        chroma_dir, collection_name = get_chroma_settings()

        try:
            count_before = vectordb._collection.count()
        except Exception:
            count_before = None

        try:
            vectordb.add_documents(docs)
            vectordb.persist()
        except ValueError as err:
            message = str(err)
            if "does not match collection dimensionality" in message:
                hint = (
                    "임베딩 차원이 기존 Chroma 컬렉션 설정과 맞지 않습니다. "
                    "새로운 임베딩 모델을 사용할 때는 `python wipe_chroma.py`로 "
                    "기존 컬렉션을 초기화하거나 `CHROMA_COLLECTION` 환경 변수를 "
                    "변경해 다른 컬렉션을 사용해야 합니다."
                )
                logger.error("Chroma dimension mismatch detected: %s", message)
                logger.warning(
                    "기존 컬렉션 %s(%s)을 삭제하고 새 임베딩으로 재시도합니다.",
                    collection_name,
                    chroma_dir,
                )
                try:
                    reset_collection()
                    vectordb = get_vectordb(emb)
                    chroma_dir, collection_name = get_chroma_settings()
                    count_before = 0
                    vectordb.add_documents(docs)
                    vectordb.persist()
                    logger.info(
                        "Chroma 컬렉션 %s 재생성 후 업로드를 완료했습니다.",
                        collection_name,
                    )
                except Exception as retry_exc:
                    logger.exception("Chroma 컬렉션 자동 초기화 재시도 실패")
                    raise HTTPException(
                        status_code=500,
                        detail=f"ingest failed: {message}. {hint}",
                    ) from retry_exc
            else:
                raise

        try:
            count_after = vectordb._collection.count()
        except Exception:
            count_after = None

        print(
            f"[INGEST] to {collection_name} @ {chroma_dir}, "
            f"file={file.filename}, chunks={len(docs)} → persist() done "
            f"(count {count_before} → {count_after})"
        )

        return {
            "ok": True,
            "file": file.filename,
            "chunks": len(docs),
            "collection": collection_name,
            "dir": chroma_dir,
            "count_before": count_before,
            "count_after": count_after,
            "view": "/admin/rag/view?limit=50&offset=0",
            "docId": doc_id,
        }

    except HTTPException as http_exc:
        logger.warning(
            "업로드 요청이 실패했습니다 status=%s detail=%s",
            http_exc.status_code,
            http_exc.detail,
        )
        raise
    except Exception as e:
        logger.exception("업로드 처리 중 예상치 못한 오류가 발생했습니다")
        raise HTTPException(status_code=500, detail=f"ingest failed: {e}")


# -----------------------------
# RAG 검색 / 생성 엔드포인트
# -----------------------------


@app.post(
    "/query/retrieve",
    response_model=RetrieveResponse,
    summary="벡터 검색 결과만 반환",
    tags=["rag"],
)
async def query_retrieve(payload: QueryRequest) -> RetrieveResponse:
    try:
        metadata_filter = {"docId": payload.doc_id} if payload.doc_id else None

        docs = _run_retrieval_with_recovery(
            question=payload.question,
            top_k=payload.top_k,
            metadata_filter=metadata_filter,
        )

        context_chunks, context_text, _has_docs = _prepare_prompt_context(docs)

        matches = [
            Match(
                reference=chunk.reference,
                chunk_index=chunk.chunk_index,
                content=chunk.body,
                preview=chunk.preview,
                source=chunk.source,
                page=chunk.page,
                metadata=chunk.metadata,
            )
            for chunk in context_chunks
        ]

        return RetrieveResponse(matches=matches, context=context_text)

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"retrieve failed: {e}")


@app.post(
    "/query/generate",
    response_model=GenerateResponse,
    summary="검색 컨텍스트로 LLM 답변 생성",
    tags=["rag"],
)
async def query_generate(payload: GenerateRequest) -> GenerateResponse:
    """
    검색된 컨텍스트를 기반으로 LLM 답변만 생성.
    (OpenAI Responses + web_search 사용)
    """
    try:
        base_context = (payload.context or "").strip()
        has_context = bool(base_context)

        answer_text = _generate_answer(
            payload.question,
            base_context,
            has_context=has_context,
        )

        return GenerateResponse(answer=answer_text)
    except HTTPException:
        raise
    except Exception as e:
        logger.exception("[QUERY_GENERATE][FATAL] q=%s err=%s", payload.question, e)
        raise HTTPException(status_code=500, detail=f"generate failed: {e}")


@app.post(
    "/query",
    response_model=QueryResponse,
    summary="벡터 검색 + LLM 기반 답변",
    tags=["rag"],
)
async def query(payload: QueryRequest) -> QueryResponse:
    try:
        metadata_filter = {"docId": payload.doc_id} if payload.doc_id else None

        docs = _run_retrieval_with_recovery(
            question=payload.question,
            top_k=payload.top_k,
            metadata_filter=metadata_filter,
        )

        context_chunks, context_text, has_docs = _prepare_prompt_context(docs)

        matches = [
            Match(
                reference=chunk.reference,
                chunk_index=chunk.chunk_index,
                content=chunk.body,
                preview=chunk.preview,
                source=chunk.source,
                page=chunk.page,
                metadata=chunk.metadata,
            )
            for chunk in context_chunks
        ]

        answer_text = _generate_answer(payload.question, context_text, has_context=has_docs)

        return QueryResponse(answer=answer_text, matches=matches)

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"query failed: {e}")


# -----------------------------
# 문서 삭제
# -----------------------------


@app.post("/documents/delete")
async def delete_documents(payload: DocDeleteRequest):
    """문서 UUID 또는 파일명을 기준으로 벡터 DB와 업로드 파일을 정리"""

    try:
        embedding_fn = get_embedding_fn()
        vectordb = get_vectordb(embedding_fn)
        collection = vectordb._collection

        if payload.doc_id:
            where: Dict[str, Any] = {"docId": payload.doc_id}
        else:
            where = {"source": payload.source}

        matched = collection.get(where=where, include=["metadatas"])
        matched_ids = matched.get("ids") or []
        metadatas = matched.get("metadatas") or []
        sources = {
            (meta or {}).get("source")
            for meta in metadatas
            if meta and (meta.get("source") is not None)
        }

        if matched_ids:
            collection.delete(where=where)
            vectordb.persist()
        else:
            logger.info(
                "삭제 요청과 일치하는 청크가 없어 delete 작업을 건너뜁니다. filter=%s",
                where,
            )

        uploads_dir = (Path(__file__).parent / "uploads").resolve()
        file_results: List[Dict[str, Any]] = []
        for file_name in sorted(s for s in sources if s):
            target_path = uploads_dir / file_name
            try:
                target_path.unlink(missing_ok=True)
                file_results.append(
                    {
                        "file": file_name,
                        "removed": not target_path.exists(),
                    }
                )
            except Exception as exc:
                file_results.append(
                    {
                        "file": file_name,
                        "removed": False,
                        "error": str(exc),
                    }
                )

        return {
            "deletedChunks": len(matched_ids),
            "deletedFiles": file_results,
        }

    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"delete failed: {exc}")


# -----------------------------
# (선택) 외부 웹 검색 API를 붙이고 싶다면
# 아래에 helper를 추가해서 `_build_openai_prompt`에 합치거나
# `/query/generate`에서 context 병합을 해도 된다.
# 현재 버전은 OpenAI의 web_search tool만 사용.
# -----------------------------


# 예시용 변수(지금은 사용하지 않음)
WEB_SEARCH_ENDPOINT = os.getenv("WEB_SEARCH_ENDPOINT")
WEB_SEARCH_API_KEY = os.getenv("WEB_SEARCH_API_KEY")


def should_use_web_search(question: str) -> bool:
    """
    필요하다면 이 함수를 사용해 외부 검색 API 여부를 판단해서
    build_web_search_context와 합쳐서 쓸 수 있음.
    (현재는 OpenAI web_search tool만 사용하므로 미사용)
    """
    q = (question or "").lower()
    time_keywords = ["현재", "지금", "요즘", "최신", "최근", "오늘", "어제", "이번 주", "이번주"]
    live_keywords = [
        "환율",
        "달러",
        "usd",
        "주가",
        "비트코인",
        "코스피",
        "나스닥",
        "금리",
        "부동산",
        "집값",
        "날씨",
        "기온",
        "뉴스",
        "이슈",
        "사건",
        "사고",
        "선거",
        "투표율",
        "경기 결과",
        "스코어",
    ]
    return any(kw in q for kw in time_keywords + live_keywords)


async def build_web_search_context(question: str, max_results: int = 5) -> str:
    """
    외부 검색엔진을 쓰고 싶을 때만 구현해서 사용.
    현재 로직에서는 호출하지 않음.
    """
    if not WEB_SEARCH_ENDPOINT or not WEB_SEARCH_API_KEY:
        return ""

    try:
        async with httpx.AsyncClient(timeout=10.0) as client:
            resp = await client.get(
                WEB_SEARCH_ENDPOINT,
                params={
                    "q": question,
                    "num": max_results,
                    "api_key": WEB_SEARCH_API_KEY,
                },
            )
            resp.raise_for_status()
            data = resp.json()
    except Exception as e:
        logger.warning("[WEB_SEARCH][ERROR] q=%s err=%s", question, e)
        return ""

    items = data.get("results") or data.get("items") or []

    lines: List[str] = []
    for idx, item in enumerate(items[:max_results], start=1):
        title = (item.get("title") or "").strip()
        url = (item.get("url") or item.get("link") or "").strip()
        snippet = (
            item.get("snippet")
            or item.get("summary")
            or item.get("description")
            or ""
        ).strip()

        if not (title or snippet):
            continue

        line = f"[검색결과 {idx}] {title}\n"
        if url:
            line += f"URL: {url}\n"
        if snippet:
            line += f"요약: {snippet}"
        lines.append(line)

    if not lines:
        return ""

    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    header = f"[웹 검색 결과]\n- 조회 시각: {now} (서버 기준)\n"
    return header + "\n\n" + "\n\n".join(lines)
